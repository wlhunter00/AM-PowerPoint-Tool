from pptx import Presentation
import six
import copy

testPresentation = 'C:\\Users\\wlhun\\OneDrive\\Documents\\Github\\AM-PowerPoint-Tool\\endTest.pptx'
bigTest = 'C:\\Users\\wlhun\\OneDrive\\Documents\\Github\\AM-PowerPoint-Tool\\bigTest.pptx'


def createPowerPointFromScratch():
    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = 'Adding a Bullet Slide'

    tf = body_shape.text_frame
    tf.text = 'Find the bullet slide layout'

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.text for first bullet'
    p.level = 1

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
    p.level = 2
    prs.save(testPresentation)


def scrapePowerPoint(presentationPath):
    prs = Presentation(presentationPath)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    print(text_runs)

# scrapePowerPoint('C:\\Users\\wlhun\\OneDrive\\Documents\\Github\\AM-PowerPoint-Tool\\test.pptx')


def delete_slide(prs, slide):
    # Make dictionary with necessary information
    id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]


def duplicate_slide(pres, index):
        template = pres.slides[index]
        try:
            blank_slide_layout = pres.slide_layouts[12]
        except:
            blank_slide_layout = pres.slide_layouts[len(pres.slide_layouts)]

        copied_slide = pres.slides.add_slide(blank_slide_layout)

        for shp in template.shapes:
            el = shp.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

        for _, value in six.iteritems(template.part.rels):
            # Make sure we don't copy a notesSlide relation as that won't exist
            if "notesSlide" not in value.reltype:
                copied_slide.part.rels.add_relationship(value.reltype,
                                                value._target,
                                                value.rId)

        return copied_slide


def learnFromScratch():
    prs = Presentation(testPresentation)
    slides = prs.slides
    print(len(slides))
    # delete_slide(prs, slides[1])
    duplicate_slide(prs, 1)
    print(len(prs.slides))
    prs.save(bigTest)


learnFromScratch()
